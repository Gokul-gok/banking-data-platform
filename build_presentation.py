from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# COLOR PALETTE
# ============================================================
NAVY        = RGBColor(0x1B, 0x31, 0x55)
NAVY2       = RGBColor(0x0F, 0x1E, 0x35)
BLUE        = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)
PALE_BG     = RGBColor(0xF0, 0xF4, 0xF9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0x6B, 0x72, 0x80)
DARK_GRAY   = RGBColor(0x37, 0x41, 0x51)
GREEN       = RGBColor(0x10, 0x7C, 0x41)
ORANGE      = RGBColor(0xC9, 0x6A, 0x10)
RED         = RGBColor(0xC0, 0x39, 0x2B)
GOLD_COLOR  = RGBColor(0xC9, 0xA0, 0x2B)
TEAL        = RGBColor(0x0D, 0x7C, 0x6E)
PURPLE      = RGBColor(0x6A, 0x3D, 0x9A)
SILVER      = RGBColor(0xA8, 0xA8, 0xA8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

HEADER_H  = Inches(1.55)
GOLD_LINE = Pt(4)
BODY_Y    = Inches(1.75)
FONT      = "Segoe UI"


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_background(slide, color=PALE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line_color=None, line_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w or Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=16, color=DARK_GRAY, bold=False,
             align=PP_ALIGN.LEFT, font=FONT, italic=False, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def add_bullets(slide, x, y, w, h, items, size=15, color=DARK_GRAY,
                line_spacing=1.4, bullet_color=None, font=FONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        text, level = (item if isinstance(item, tuple) else (item, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(6)
        p.line_spacing = line_spacing
        bullet = "▸  " if level == 0 else "  –  "
        run = p.add_run()
        run.text = f"{bullet}{text}"
        run.font.size = Pt(size - (2 if level else 0))
        run.font.color.rgb = bullet_color or color
        run.font.name = font
    return box


def slide_header(slide, title, subtitle=None):
    """Dark navy header with gold accent line"""
    # Navy background band
    add_rect(slide, 0, 0, prs.slide_width, HEADER_H, NAVY)
    # Subtle inner gradient effect (lighter strip at top)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.06), RGBColor(0x28, 0x48, 0x72))
    # Gold accent line below header
    add_rect(slide, 0, HEADER_H, prs.slide_width, GOLD_LINE, GOLD_COLOR)
    # Left accent stripe
    add_rect(slide, 0, 0, Inches(0.07), prs.slide_height, GOLD_COLOR)
    # Title
    add_text(slide, Inches(0.55), Inches(0.18), Inches(12.2), Inches(0.85),
             title, size=28, color=WHITE, bold=True, font=FONT)
    # Subtitle
    if subtitle:
        add_text(slide, Inches(0.55), Inches(0.98), Inches(12.2), Inches(0.42),
                 subtitle, size=13, color=LIGHT_BLUE, italic=True, font=FONT)


def add_footer(slide, page_num):
    add_rect(slide, 0, Inches(7.18), prs.slide_width, Pt(2), GOLD_COLOR)
    add_text(slide, Inches(0.55), Inches(7.22), Inches(9), Inches(0.25),
             "Banking Data Platform  ·  Batch & Real-Time Streaming Architecture",
             size=9, color=GRAY, italic=True)
    add_text(slide, Inches(12.5), Inches(7.22), Inches(0.5), Inches(0.25),
             str(page_num), size=9, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, title_bg=NAVY, title_color=WHITE,
             body_bg=WHITE, body_color=DARK_GRAY, title_size=14, body_size=12):
    """Card with colored header band + white body"""
    th = Inches(0.42)
    # Title band
    add_rect(slide, x, y, w, th, title_bg)
    add_text(slide, x+Inches(0.15), y, w-Inches(0.3), th,
             title, size=title_size, color=title_color, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    # Body band
    add_rect(slide, x, y+th, w, h-th, body_bg, line_color=RGBColor(0xCC,0xD9,0xEA), line_w=Pt(0.75))
    add_text(slide, x+Inches(0.15), y+th+Inches(0.06), w-Inches(0.3), h-th-Inches(0.1),
             body, size=body_size, color=body_color)


def add_stat_card(slide, x, y, w, h, value, label, val_color=NAVY, bg=WHITE):
    """KPI stat card"""
    add_rect(slide, x, y, w, h, bg, line_color=LIGHT_BLUE, line_w=Pt(1.5))
    add_rect(slide, x, y, Inches(0.07), h, BLUE)
    add_text(slide, x+Inches(0.18), y+Inches(0.1), w-Inches(0.28), Inches(0.55),
             value, size=26, color=val_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x+Inches(0.18), y+Inches(0.62), w-Inches(0.28), Inches(0.28),
             label, size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1: TITLE
# ============================================================
s = add_slide()
set_background(s, NAVY2)
# Background design elements
add_rect(s, 0, 0, Inches(0.12), prs.slide_height, GOLD_COLOR)
add_rect(s, prs.slide_width-Inches(0.12), 0, Inches(0.12), prs.slide_height, GOLD_COLOR)
add_rect(s, 0, Inches(3.52), prs.slide_width, Pt(3), GOLD_COLOR)
add_rect(s, 0, Inches(3.56), prs.slide_width, Pt(1), RGBColor(0xFF,0xFF,0xFF))

# Title block
add_text(s, Inches(1), Inches(1.6), Inches(11.3), Inches(1.1),
         "Banking Data Platform", size=46, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, font=FONT)
add_text(s, Inches(1), Inches(2.65), Inches(11.3), Inches(0.6),
         "End-to-End Lakehouse Architecture", size=22, color=LIGHT_BLUE,
         align=PP_ALIGN.CENTER, italic=True)

# Tech chips
chips = ["Azure Data Lake", "Databricks", "Event Hubs", "Synapse Analytics", "Power BI", "ADF", "Logic Apps"]
chip_x = Inches(1.0)
chip_y = Inches(4.3)
chip_w = Inches(1.6)
for chip in chips:
    add_rect(s, chip_x, chip_y, chip_w, Inches(0.38), RGBColor(0x28, 0x48, 0x72))
    add_text(s, chip_x, chip_y, chip_w, Inches(0.38),
             chip, size=10, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    chip_x += chip_w + Inches(0.1)

add_text(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.38),
         "Prepared by  Gokul Sankar", size=13, color=GRAY,
         align=PP_ALIGN.CENTER, italic=True)


# ============================================================
# SLIDE 2: AGENDA
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Agenda")

agenda_left = [
    "Project Overview & Business Objective",
    "Azure Services Used",
    "Medallion Architecture (Bronze / Silver / Gold)",
    "Batch Pipeline — Ingestion to Gold",
    "Gold Layer — Star Schema Model",
    "Streaming Pipeline — Event Hub to Gold",
]
agenda_right = [
    "Streaming Gold Tables",
    "ADF Orchestration",
    "Monitoring & Email Notifications",
    "Power BI — Batch Analytics Dashboard",
    "Power BI — Real-Time Streaming Dashboard",
    "Key Business Insights  ·  Challenges  ·  Future",
]

# Two-column numbered agenda
for i, item in enumerate(agenda_left, 1):
    y = BODY_Y + (i-1) * Inches(0.68)
    add_rect(s, Inches(0.55), y, Inches(0.5), Inches(0.42), NAVY)
    add_text(s, Inches(0.55), y, Inches(0.5), Inches(0.42), str(i),
             size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.18), y+Inches(0.04), Inches(5.2), Inches(0.42),
             item, size=14, color=DARK_GRAY)

for i, item in enumerate(agenda_right, 7):
    y = BODY_Y + (i-7) * Inches(0.68)
    add_rect(s, Inches(7.0), y, Inches(0.5), Inches(0.42), BLUE)
    add_text(s, Inches(7.0), y, Inches(0.5), Inches(0.42), str(i),
             size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.63), y+Inches(0.04), Inches(5.2), Inches(0.42),
             item, size=14, color=DARK_GRAY)

add_rect(s, Inches(6.7), BODY_Y, Pt(2), Inches(4.2), LIGHT_BLUE)
add_footer(s, 2)


# ============================================================
# SLIDE 3: PROJECT OVERVIEW
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Project Overview", "What this project builds and delivers")

add_card(s, Inches(0.55), BODY_Y, Inches(5.9), Inches(2.3),
         "Business Objective",
         "Build a full lakehouse platform for a retail bank — ingesting customer, account, "
         "card, and loan data, plus real-time transaction, fraud, and login event streams — "
         "to power analytics on churn, gap analysis, city performance, and live fraud monitoring.",
         title_bg=NAVY, title_color=WHITE, body_color=DARK_GRAY)

add_card(s, Inches(6.7), BODY_Y, Inches(6.05), Inches(2.3),
         "Scope",
         "Two parallel pipelines feeding the same Gold and Power BI layer:\n"
         "  Batch: historical CSV datasets (customers, accounts, cards, loans)\n"
         "  Streaming: live transactions, fraud alerts, login events via Event Hub",
         title_bg=BLUE, title_color=WHITE, body_color=DARK_GRAY)

add_card(s, Inches(0.55), BODY_Y+Inches(2.5), Inches(12.2), Inches(2.3),
         "Key Outcomes",
         "Two ADF-orchestrated pipelines (batch + streaming) with scheduled triggers  ·  "
         "Star-schema Gold model (dim/fact) for flexible analytics  ·  "
         "9 streaming Gold tables for near real-time monitoring  ·  "
         "Two Power BI dashboards: Historical Analytics & Real-Time Operations  ·  "
         "Automated email alerts on pipeline success/failure via Logic Apps",
         title_bg=GOLD_COLOR, title_color=NAVY, body_color=DARK_GRAY, body_size=13)
add_footer(s, 3)


# ============================================================
# SLIDE 4: AZURE SERVICES
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Azure Services Used", "Eight Azure services working together end-to-end")

services = [
    (BLUE,       "Azure Data Lake Storage Gen2",       "Bronze / Silver / Gold containers — single source of truth across all pipeline layers"),
    (NAVY,       "Azure Key Vault",                     "Secure storage for storage account keys and Event Hub connection strings"),
    (RGBColor(0x6A,0x3D,0x9A), "Azure Databricks",    "Spark-based ETL — Bronze ingestion, Silver cleansing, Gold aggregation (batch + streaming)"),
    (ORANGE,     "Azure Event Hubs",                   "Real-time ingestion of transactions, fraud alerts, and login events via Kafka protocol"),
    (TEAL,       "Azure Data Factory",                 "Orchestrates all pipelines — scheduling, dependency chains, failure handling"),
    (GREEN,      "Azure Synapse Analytics",             "Serverless SQL views over Gold Delta files — query layer for Power BI"),
    (RED,        "Power BI",                            "Two dashboards — historical batch analytics and real-time streaming monitor"),
    (RGBColor(0x80,0x40,0x00), "Azure Logic Apps",    "Email notifications on pipeline success / failure via SMTP (Gmail App Password)"),
]

y = BODY_Y
for color, name, desc in services:
    add_rect(s, Inches(0.55), y, Inches(0.08), Inches(0.58), color)
    add_rect(s, Inches(0.63), y, Inches(12.1), Inches(0.58), WHITE,
             line_color=RGBColor(0xCC,0xD9,0xEA), line_w=Pt(0.5))
    add_text(s, Inches(0.82), y+Inches(0.05), Inches(3.3), Inches(0.48),
             name, size=13, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(4.2), y+Inches(0.15), Pt(1.5), Inches(0.28), LIGHT_BLUE)
    add_text(s, Inches(4.4), y+Inches(0.05), Inches(8.1), Inches(0.48),
             desc, size=12, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.63)
add_footer(s, 4)


# ============================================================
# SLIDE 5: MEDALLION ARCHITECTURE
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Medallion Architecture", "Bronze  →  Silver  →  Gold — progressive data refinement")

layers = [
    ("RAW",    "Source Data",           RGBColor(0x55,0x65,0x78), RGBColor(0x37,0x41,0x51),
     "Unmodified source data\nADLS raw container\nEvent Hub (streaming)"),
    ("BRONZE", "As-Is Ingestion",       RGBColor(0xA0,0x52,0x22), RGBColor(0x7A,0x38,0x10),
     "CSV → Delta format\nAdd tracking columns:\n_source_file, _ingested_at"),
    ("SILVER", "Cleaned & Conformed",   SILVER,                   RGBColor(0x77,0x77,0x77),
     "Type casting (dates, numerics)\nStandardise text (initcap)\nDrop duplicates"),
    ("GOLD",   "Business-Ready",        GOLD_COLOR,               RGBColor(0x9A,0x72,0x10),
     "Star schema (batch)\nAggregated tables (streaming)\nSynapse views → Power BI"),
]

x = Inches(0.5)
w = Inches(2.9)
gap = Inches(0.25)
box_y = BODY_Y + Inches(0.2)
box_h = Inches(4.3)

for i, (tag, title, header_color, band_color, desc) in enumerate(layers):
    # Main box
    add_rect(s, x, box_y, w, box_h, WHITE, line_color=header_color, line_w=Pt(1.5))
    # Colored header band
    add_rect(s, x, box_y, w, Inches(0.9), header_color)
    # Tag label
    add_text(s, x+Inches(0.1), box_y+Inches(0.05), w-Inches(0.2), Inches(0.48),
             tag, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Title below tag
    add_text(s, x+Inches(0.1), box_y+Inches(0.52), w-Inches(0.2), Inches(0.34),
             title, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    # Description in body
    add_text(s, x+Inches(0.15), box_y+Inches(1.08), w-Inches(0.3), Inches(3.0),
             desc, size=13, color=DARK_GRAY)
    # Arrow
    if i < len(layers) - 1:
        add_text(s, x+w+Inches(0.02), box_y+Inches(1.8), gap, Inches(0.5),
                 "→", size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    x += w + gap

add_footer(s, 5)


# ============================================================
# SLIDE 6: BATCH PIPELINE FLOW
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Batch Pipeline", "CSV datasets  →  Bronze  →  Silver  →  Gold (Star Schema)")

steps = [
    (NAVY,       "01  Bronze Ingestion",
     "Reads accounts, cards, customers, loans CSVs from ADLS raw container.\n"
     "Adds _source_file and _ingested_at tracking columns.\n"
     "Writes as Delta format — no data changes, schema preserved."),
    (BLUE,       "02  Silver Transformation",
     "Casts date columns (dd-MM-yyyy  →  DATE type).\n"
     "Casts numeric types (balance  →  double, credit_score  →  int).\n"
     "Trims whitespace, standardises text casing with initcap."),
    (GOLD_COLOR, "03  Gold — Star Schema",
     "Dimensions: dim_customer, dim_date, dim_account_type, dim_risk\n"
     "Facts: fact_account, fact_loan, fact_card\n"
     "All layers written to ADLS as Delta, read by Synapse views."),
]

x = Inches(0.55)
w = Inches(3.9)
y_box = BODY_Y

for i, (color, title, desc) in enumerate(steps):
    add_rect(s, x, y_box, w, Inches(0.58), color)
    add_text(s, x+Inches(0.15), y_box, w-Inches(0.3), Inches(0.58),
             title, size=16, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, y_box+Inches(0.58), w, Inches(2.4), WHITE,
             line_color=color, line_w=Pt(1.5))
    add_text(s, x+Inches(0.18), y_box+Inches(0.72), w-Inches(0.36), Inches(2.1),
             desc, size=13, color=DARK_GRAY)
    if i < 2:
        add_text(s, x+w+Inches(0.02), y_box+Inches(1.3), Inches(0.22), Inches(0.5),
                 "→", size=24, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    x += w + Inches(0.23)

# Business outputs banner
add_rect(s, Inches(0.55), BODY_Y+Inches(3.25), Inches(12.2), Pt(3), GOLD_COLOR)
add_rect(s, Inches(0.55), BODY_Y+Inches(3.32), Inches(12.2), Inches(1.55),
         RGBColor(0xFD,0xF6,0xE3), line_color=GOLD_COLOR, line_w=Pt(1))
add_text(s, Inches(0.75), BODY_Y+Inches(3.38), Inches(2.0), Inches(0.35),
         "Business Outputs", size=13, color=GOLD_COLOR, bold=True)
add_text(s, Inches(0.75), BODY_Y+Inches(3.72), Inches(11.8), Inches(1.0),
         "Customer Churn  ·  Signup-to-Account Gap Analysis  ·  City Metrics  ·  "
         "Monthly New Account & Loan Trend  ·  Loan Risk Analysis  ·  Customer 360",
         size=13, color=DARK_GRAY)
add_footer(s, 6)


# ============================================================
# SLIDE 7: GOLD LAYER -- TABLE LIST
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Gold Layer -- Batch Tables", "Business-ready tables written to ADLS and exposed via Synapse views")

dim_fact_tables = [
    (GOLD_COLOR, "Dimensions",  ["dim_customer", "dim_date", "dim_account_type", "dim_risk"]),
    (NAVY,       "Fact Tables", ["fact_account", "fact_loan", "fact_card"]),
]

# Left section -- Dimensions
add_rect(s, Inches(0.55), BODY_Y, Inches(5.9), Inches(0.5), GOLD_COLOR)
add_text(s, Inches(0.55), BODY_Y, Inches(5.9), Inches(0.5),
         "Dimensions", size=16, color=NAVY, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

dim_tables = [
    ("dim_customer",     "Stores customer details: name, city, credit score, credit tier"),
    ("dim_date",         "Full calendar table: day, month, quarter, year, weekend flag"),
    ("dim_account_type", "Account type lookup: Savings, Checking, Business"),
    ("dim_risk",         "Risk band lookup: Low Risk, Medium Risk, High Risk"),
]

y = BODY_Y + Inches(0.6)
for table, desc in dim_tables:
    add_rect(s, Inches(0.55), y, Inches(0.08), Inches(0.72), GOLD_COLOR)
    add_rect(s, Inches(0.63), y, Inches(5.82), Inches(0.72), WHITE,
             line_color=GOLD_COLOR, line_w=Pt(0.75))
    add_text(s, Inches(0.85), y+Inches(0.05), Inches(5.3), Inches(0.3),
             table, size=14, color=NAVY, bold=True)
    add_text(s, Inches(0.85), y+Inches(0.36), Inches(5.3), Inches(0.3),
             desc, size=11, color=GRAY)
    y += Inches(0.82)

# Right section -- Facts
add_rect(s, Inches(6.9), BODY_Y, Inches(5.9), Inches(0.5), NAVY)
add_text(s, Inches(6.9), BODY_Y, Inches(5.9), Inches(0.5),
         "Fact Tables", size=16, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

fact_tables = [
    ("fact_account", "One row per account: balance, status, churn, gap analysis fields"),
    ("fact_loan",    "One row per loan: amount, interest rate, risk category"),
    ("fact_card",    "One row per card: card type, linked to account and customer"),
]

y = BODY_Y + Inches(0.6)
for table, desc in fact_tables:
    add_rect(s, Inches(6.9), y, Inches(0.08), Inches(0.72), NAVY)
    add_rect(s, Inches(6.98), y, Inches(5.82), Inches(0.72), WHITE,
             line_color=NAVY, line_w=Pt(0.75))
    add_text(s, Inches(7.2), y+Inches(0.05), Inches(5.3), Inches(0.3),
             table, size=14, color=NAVY, bold=True)
    add_text(s, Inches(7.2), y+Inches(0.36), Inches(5.3), Inches(0.3),
             desc, size=11, color=GRAY)
    y += Inches(0.82)

# Divider
add_rect(s, Inches(6.7), BODY_Y, Pt(2), Inches(4.5), LIGHT_BLUE)

add_footer(s, 7)

# SLIDE 8: STREAMING PIPELINE FLOW
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Streaming Pipeline", "Event Producer  →  Event Hub  →  Bronze  →  Silver  →  Gold (every 15 min)")

stages = [
    (ORANGE,     "Event Producer\n(Python, local)",
     "3 generators in parallel threads\n(transactions, fraud, login events)"),
    (RGBColor(0x6A,0x3D,0x9A), "Event Hub",
     "Single hub: banking-transactions\nevent_type field routes records"),
    (NAVY,       "Bronze\n(Databricks Streaming)",
     "spark.readStream from Event Hub\n3 continuous Delta writers"),
    (BLUE,       "Silver\n(Databricks Streaming)",
     "Cleansing + derived columns\n(amount_category, session_category)"),
    (GOLD_COLOR, "Gold\n(Databricks Batch, 15-min trigger)",
     "9 aggregated tables\nviews created in Synapse\nto connect to Power BI"),
]

x = Inches(0.5)
w = Inches(2.3)
bx_y = BODY_Y + Inches(0.3)
bx_h = Inches(3.5)

for i, (color, title, desc) in enumerate(stages):
    add_rect(s, x, bx_y, w, Inches(0.88), color)
    tf_box = s.shapes.add_textbox(x, bx_y, w, Inches(0.88))
    tf = tf_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(11.5)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = FONT

    add_rect(s, x, bx_y+Inches(0.88), w, bx_h-Inches(0.88), WHITE,
             line_color=color, line_w=Pt(1.5))
    add_text(s, x+Inches(0.12), bx_y+Inches(1.0), w-Inches(0.24),
             bx_h-Inches(1.1), desc, size=12, color=DARK_GRAY)

    if i < len(stages) - 1:
        add_text(s, x+w+Inches(0.01), bx_y+Inches(1.5), Inches(0.25), Inches(0.5),
                 "→", size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    x += w + Inches(0.14)

add_card(s, Inches(0.55), BODY_Y+Inches(3.95), Inches(12.2), Inches(1.5),
         "Why Gold Runs in a Separate Pipeline",
         "Bronze and Silver are continuous streams that never finish — they run in their own ADF pipeline triggered once. "
         "Gold re-aggregates Silver on a 15-minute schedule (streaming_gold_pipeline) since it is a "
         "batch-style summarisation, not a continuous stream.",
         title_bg=NAVY, body_color=DARK_GRAY)
add_footer(s, 8)


# ============================================================
# SLIDE 9: STREAMING GOLD TABLES
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Streaming Gold Tables", "9 aggregated tables for real-time dashboards")

cols = [
    (BLUE,   "Transactions",
     ["stream_txn_hourly — volume & amount by hour",
      "stream_txn_city — activity by city",
      "stream_txn_type — type x channel breakdown"]),
    (RED,    "Fraud Alerts",
     ["stream_fraud_summary — by type/severity/action",
      "stream_fraud_city — alerts by city",
      "stream_fraud_hourly — trend & blocked count"]),
    (GREEN,  "Login Events",
     ["stream_login_summary — success/fail by method/device",
      "stream_login_failures — failure reasons by city",
      "stream_login_hourly — login activity trend"]),
]

x = Inches(0.55)
w = Inches(3.95)
for title, items, color in [(c[1], c[2], c[0]) for c in cols]:
    add_rect(s, x, BODY_Y, w, Inches(0.6), color)
    add_text(s, x, BODY_Y, w, Inches(0.6), title, size=17, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, BODY_Y+Inches(0.6), w, Inches(3.8), WHITE,
             line_color=color, line_w=Pt(1.5))
    y_item = BODY_Y + Inches(0.85)
    for item in items:
        add_rect(s, x+Inches(0.2), y_item+Inches(0.04), Inches(0.22), Inches(0.22), color)
        add_text(s, x+Inches(0.55), y_item, w-Inches(0.7), Inches(0.3),
                 item, size=13, color=DARK_GRAY)
        y_item += Inches(0.95)
    x += w + Inches(0.28)

add_footer(s, 9)


# ============================================================
# SLIDE 10: ADF ORCHESTRATION
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Orchestration with Azure Data Factory", "Three pipelines covering batch, streaming ingestion, and streaming gold")

pipelines = [
    (NAVY,       "banking_batch_pipeline",    "Manual / Scheduled",
     "run_bronze  →  run_silver  →  run_gold",
     "Success & failure email via Logic App on every step"),
    (ORANGE,     "streaming_pipeline",        "Run once (continuous)",
     "streaming_bronze   ‖   streaming_silver (parallel, never stop)",
     "Failure-only alert — continuous streams have no natural completion"),
    (GOLD_COLOR, "streaming_gold_pipeline",   "Every 15 minutes",
     "streaming_gold  (re-aggregates Silver into 9 Gold tables)",
     "Success & failure email via Logic App on every run"),
]

y = BODY_Y
for color, name, trigger, flow, notif in pipelines:
    add_rect(s, Inches(0.55), y, Inches(0.08), Inches(1.65), color)
    add_rect(s, Inches(0.63), y, Inches(12.1), Inches(1.65), WHITE,
             line_color=color, line_w=Pt(0.75))
    add_text(s, Inches(0.82), y+Inches(0.1), Inches(5.5), Inches(0.42),
             name, size=16, color=NAVY, bold=True)
    add_rect(s, Inches(0.82), y+Inches(0.6), Inches(2.0), Inches(0.35), color)
    add_text(s, Inches(0.82), y+Inches(0.6), Inches(2.0), Inches(0.35),
             trigger, size=11, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.02), y+Inches(0.6), Inches(9.5), Inches(0.35),
             flow, size=13, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.82), y+Inches(1.1), Inches(11.7), Inches(0.42),
             "  ✉  " + notif, size=12, color=GRAY, italic=True)
    y += Inches(1.82)
add_footer(s, 10)


# ============================================================
# SLIDE 11: MONITORING / NOTIFICATIONS
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Monitoring & Notifications", "Azure Logic Apps + SMTP email alerts on every pipeline run")

steps = [
    ("1", NAVY,  "ADF Web Activity",
     "Calls the Logic App HTTP trigger at the end of each pipeline activity (green arrow = success path, red arrow = failure path)"),
    ("2", BLUE,  "Logic App Workflow",
     "Receives PipelineName, Status, and Time as a JSON payload via HTTP POST"),
    ("3", GREEN, "Send Email (SMTP)",
     "Gmail App Password + smtp.gmail.com:465 delivers the notification to gokulgokz14@gmail.com"),
    ("4", GOLD_COLOR, "Coverage",
     "banking_batch_pipeline (every step)  ·  streaming_gold_pipeline (every run)  ·  streaming_pipeline (failure only)"),
]

y = BODY_Y
for num, color, title, desc in steps:
    add_rect(s, Inches(0.55), y, Inches(0.6), Inches(0.88), color)
    add_text(s, Inches(0.55), y, Inches(0.6), Inches(0.88), num,
             size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(1.15), y, Inches(11.6), Inches(0.88), WHITE,
             line_color=color, line_w=Pt(0.75))
    add_text(s, Inches(1.35), y+Inches(0.06), Inches(2.5), Inches(0.35),
             title, size=14, color=color, bold=True)
    add_text(s, Inches(1.35), y+Inches(0.44), Inches(11.1), Inches(0.35),
             desc, size=12, color=DARK_GRAY)
    y += Inches(1.0)

add_card(s, Inches(0.55), y+Inches(0.1), Inches(12.2), Inches(1.0),
         "Sample JSON Payload",
         '{ "PipelineName": "streaming_gold_pipeline",  "Status": "Succeeded",  "Time": "2026-07-01T10:15:00Z" }',
         title_bg=NAVY, body_bg=RGBColor(0x1E,0x1E,0x2E), body_color=LIGHT_BLUE, body_size=12)
add_footer(s, 11)


# ============================================================
# SLIDE 12: POWER BI - BATCH DASHBOARD
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Power BI — Batch Analytics Dashboard", "5 report pages built on the Gold star schema")

pages = [
    (NAVY,       "Overview",        "5 KPI cards (customers, accounts, total balance, loans, loan amount), donut by account type, risk distribution bar chart, slicers for city / account type / year"),
    (BLUE,       "Customer Churn",  "Churn status donut (Active / Inactive / At Risk / Churned), churn by city (Top 10), churn by credit tier — severity slicer for drill-down"),
    (GOLD_COLOR, "Gap Analysis",    "Avg days from signup to first account opening, gap category breakdown bar chart, by city (Top 10) and by credit tier — filtered to is_first_account = True"),
    (GREEN,      "City Metrics",    "Top 10 cities by customers, avg balance, avg credit score, total loan exposure — city dropdown slicer for cross-filtering"),
    (ORANGE,     "Monthly Trend",   "New accounts by month / year with account-type breakdown, new loans by year, seasonal column chart by month — year slicer synced across all pages"),
]

y = BODY_Y
for i, (color, title, desc) in enumerate(pages, 1):
    add_rect(s, Inches(0.55), y, Inches(0.08), Inches(0.85), color)
    add_rect(s, Inches(0.63), y, Inches(12.1), Inches(0.85), WHITE,
             line_color=color, line_w=Pt(0.5))
    add_rect(s, Inches(0.63), y, Inches(2.3), Inches(0.85), RGBColor(0xF5,0xF7,0xFA))
    add_text(s, Inches(0.78), y, Inches(2.1), Inches(0.85),
             f"Page {i}  —  {title}", size=13, color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.1), y+Inches(0.08), Inches(9.4), Inches(0.72),
             desc, size=12, color=DARK_GRAY)
    y += Inches(0.98)
add_footer(s, 12)


# ============================================================
# SLIDE 13: POWER BI - STREAMING DASHBOARD
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Power BI — Real-Time Streaming Dashboard", "3 pages built on Gold tables via Synapse views")

pages2 = [
    (BLUE,  "Transaction Monitor",
     ["Total transactions / amount / avg amount / failed transactions (KPI cards)",
      "Transaction volume over time (line chart by hour)",
      "Completed vs Failed vs Pending by day (stacked column)",
      "Transactions by city (Top 10 bar), by type (donut), by channel (bar)"]),
    (RED,   "Fraud Monitor",
     ["Total alerts / amount flagged / confirmed frauds / avg risk score (KPI cards)",
      "Alerts by severity (donut) and by fraud type (bar)",
      "Action taken breakdown (Blocked / Flagged / Allowed)",
      "Fraud trend over time (line), fraud hotspot cities (bar), severity slicer"]),
    (GREEN, "Login Analytics",
     ["Total / successful / failed logins / failure rate % (KPI cards)",
      "Login success vs failure (donut), logins by method and device (bars)",
      "Failure reasons (bar), failed logins by city (bar)",
      "Login activity over time (line), device slicer for drill-down"]),
]

x = Inches(0.55)
w = Inches(3.95)
for title, items, color in [(c[1], c[2], c[0]) for c in pages2]:
    add_rect(s, x, BODY_Y, w, Inches(0.55), color)
    add_text(s, x, BODY_Y, w, Inches(0.55), title, size=15, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, BODY_Y+Inches(0.55), w, Inches(3.7), WHITE,
             line_color=color, line_w=Pt(1.5))
    y_it = BODY_Y + Inches(0.8)
    for item in items:
        add_rect(s, x+Inches(0.2), y_it+Inches(0.05), Inches(0.16), Inches(0.16), color)
        add_text(s, x+Inches(0.5), y_it, w-Inches(0.65), Inches(0.26),
                 item, size=11.5, color=DARK_GRAY)
        y_it += Inches(0.75)
    x += w + Inches(0.28)

add_card(s, Inches(0.55), BODY_Y+Inches(4.55), Inches(12.2), Inches(0.8),
         "Refresh Cadence",
         "Gold tables refresh every 15 minutes via ADF trigger  ·  Synapse views auto-read latest Delta files  ·  Power BI: manual refresh in Desktop, scheduled refresh in Power BI Service",
         title_bg=NAVY, body_color=DARK_GRAY)
add_footer(s, 13)


# ============================================================
# SLIDE 13b: PIPELINE EMAIL NOTIFICATIONS
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Pipeline Email Notifications", "Logic Apps sends success and failure emails via SMTP (Gmail)")

import os as _os

# Success card
add_rect(s, Inches(0.55), BODY_Y, Inches(5.9), Inches(0.45), GREEN)
add_text(s, Inches(0.55), BODY_Y, Inches(5.9), Inches(0.45),
         "Success Notification", size=15, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

success_img = _os.path.join(_os.path.dirname(_os.path.abspath("build_presentation.py")), "success.png")
s.shapes.add_picture(success_img, Inches(0.55), BODY_Y+Inches(0.55), Inches(5.9), Inches(4.0))

# Failure card
add_rect(s, Inches(6.85), BODY_Y, Inches(5.9), Inches(0.45), RED)
add_text(s, Inches(6.85), BODY_Y, Inches(5.9), Inches(0.45),
         "Failure Notification", size=15, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

failure_img = _os.path.join(_os.path.dirname(_os.path.abspath("build_presentation.py")), "failure.png")
s.shapes.add_picture(failure_img, Inches(6.85), BODY_Y+Inches(0.55), Inches(5.9), Inches(4.0))

# Divider
add_rect(s, Inches(6.65), BODY_Y, Pt(2), Inches(4.7), LIGHT_BLUE)

# Bottom note
add_card(s, Inches(0.55), BODY_Y+Inches(4.75), Inches(12.2), Inches(0.7),
         "How it works",
         "ADF Web Activity (green arrow = success, red arrow = failure) calls Logic App HTTP trigger  "
         "→  Logic App sends email with Pipeline Name, Status and Time via SMTP",
         title_bg=NAVY, body_color=DARK_GRAY)

add_footer(s, 14)


# ============================================================
# SLIDE 14: KEY BUSINESS INSIGHTS
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Key Business Insights Delivered", "Six actionable insights powered by this platform")

insights = [
    (NAVY,       "Customer Retention",
     "Churn segmentation (Active / Inactive / At Risk / Churned) flags accounts going dormant before they close — enabling proactive outreach before the customer is lost."),
    (BLUE,       "Onboarding Efficiency",
     "Gap analysis measures how long customers wait between signup and opening their first product — a direct indicator of onboarding friction in the digital journey."),
    (GOLD_COLOR, "Regional Strategy",
     "City-level metrics (avg balance, avg credit score, loan exposure) highlight which regions carry the most value and risk — informing branch and marketing investment decisions."),
    (GREEN,      "Seasonality Planning",
     "Monthly trend analysis reveals peak months for account openings and loan originations, helping the bank plan staffing, promotions, and system capacity in advance."),
    (RED,        "Live Fraud Defense",
     "Real-time fraud dashboard surfaces severity, fraud type, and city hotspots within minutes of detection — far faster than traditional end-of-day batch reporting."),
    (ORANGE,     "Digital Channel Health",
     "Login analytics expose failure rates by device and method, pinpointing friction in the digital banking experience before it drives customers to competitors."),
]

y = BODY_Y
for i, (color, title, desc) in enumerate(insights):
    col = 0 if i % 2 == 0 else 1
    row = i // 2
    x = Inches(0.55) if col == 0 else Inches(6.95)
    ry = BODY_Y + row * Inches(1.72)
    add_rect(s, x, ry, Inches(0.08), Inches(1.5), color)
    add_rect(s, x+Inches(0.08), ry, Inches(5.9), Inches(1.5), WHITE,
             line_color=color, line_w=Pt(0.5))
    add_text(s, x+Inches(0.25), ry+Inches(0.1), Inches(5.5), Inches(0.38),
             title, size=14, color=color, bold=True)
    add_text(s, x+Inches(0.25), ry+Inches(0.5), Inches(5.5), Inches(0.9),
             desc, size=11.5, color=DARK_GRAY)

add_footer(s, 14)


# ============================================================
# SLIDE 15: CHALLENGES & RESOLUTIONS
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Challenges & Resolutions", "Real issues encountered and how they were resolved")

rows = [
    ("Delta column-mapping options unreliable",
     "Sanitized column names at source — replaced special characters with underscores"),
    ("Azure VM core quota exceeded on cluster creation",
     "Switched to Single Node cluster (DS3_v2, 0 workers) to stay within the approved 10-core quota"),
    ("Kafka connector blocked on Unity Catalog shared clusters",
     "Used Kafka protocol built into Databricks runtime — no Maven library needed"),
    ("Kafka SASL login module not found",
     "Used the shaded class name: kafkashaded.org.apache.kafka...PlainLoginModule"),
    ("Synapse directory cannot be listed error",
     "Diagnosed as expired SAS token — fixed via ALTER DATABASE SCOPED CREDENTIAL"),
    ("Office 365 and Gmail OAuth connectors unavailable",
     "Used generic SMTP connector with a Gmail App Password — works in all regions"),
]

y = BODY_Y
for issue, fix in rows:
    # Issue box
    add_rect(s, Inches(0.55), y, Inches(0.08), Inches(0.72), RED)
    add_rect(s, Inches(0.63), y, Inches(5.6), Inches(0.72),
             RGBColor(0xFD,0xF0,0xEE), line_color=RED, line_w=Pt(0.5))
    add_text(s, Inches(0.78), y+Inches(0.05), Inches(5.25), Inches(0.62),
             issue, size=12, color=RED, anchor=MSO_ANCHOR.MIDDLE)
    # Arrow
    add_text(s, Inches(6.4), y+Inches(0.15), Inches(0.4), Inches(0.42),
             "→", size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    # Fix box
    add_rect(s, Inches(6.9), y, Inches(0.08), Inches(0.72), GREEN)
    add_rect(s, Inches(6.98), y, Inches(5.75), Inches(0.72),
             RGBColor(0xED,0xF7,0xF0), line_color=GREEN, line_w=Pt(0.5))
    add_text(s, Inches(7.13), y+Inches(0.05), Inches(5.4), Inches(0.62),
             fix, size=12, color=GREEN, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.82)
add_footer(s, 16)


# ============================================================
# SLIDE 16: FUTURE ENHANCEMENTS
# ============================================================
s = add_slide()
set_background(s, PALE_BG)
slide_header(s, "Future Enhancements", "Production-grade improvements for scale and governance")

future = [
    (NAVY,       "Retry Policies",       "Add retry count (3 attempts, 30s interval) on all ADF notebook activities for transient-failure resilience"),
    (BLUE,       "Master Pipeline",      "Single orchestrator pipeline to trigger batch + streaming + gold from one place with one click"),
    (GOLD_COLOR, "Parameterisation",     "Pipeline parameters (environment, table list) for clean dev / test / prod separation without code changes"),
    (GREEN,      "Power BI Service",     "Publish dashboards to Power BI Service with scheduled refresh (15–30 min) for true near-real-time analytics"),
    (ORANGE,     "Sub-minute Streaming", "Switch streaming Gold to Structured Streaming foreachBatch for sub-minute latency instead of 15-min batch schedule"),
    (PURPLE,     "Unity Catalog",        "Implement Unity Catalog for column-level lineage, access control, and data discovery across all three layers"),
]

y = BODY_Y
for i, (color, title, desc) in enumerate(future):
    col = i % 2
    row = i // 2
    x = Inches(0.55) if col == 0 else Inches(6.95)
    ry = BODY_Y + row * Inches(1.62)
    add_rect(s, x, ry, Inches(5.9), Inches(0.42), color)
    add_text(s, x+Inches(0.15), ry, Inches(5.6), Inches(0.42),
             title, size=14, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, x, ry+Inches(0.42), Inches(5.9), Inches(1.05), WHITE,
             line_color=color, line_w=Pt(1))
    add_text(s, x+Inches(0.15), ry+Inches(0.52), Inches(5.6), Inches(0.88),
             desc, size=12, color=DARK_GRAY)

add_footer(s, 17)


# ============================================================
# SLIDE 17: THANK YOU
# ============================================================
s = add_slide()
set_background(s, NAVY2)
add_rect(s, 0, 0, Inches(0.12), prs.slide_height, GOLD_COLOR)
add_rect(s, prs.slide_width-Inches(0.12), 0, Inches(0.12), prs.slide_height, GOLD_COLOR)
add_rect(s, 0, Inches(3.58), prs.slide_width, Pt(3), GOLD_COLOR)

add_text(s, Inches(1), Inches(2.4), Inches(11.3), Inches(1.1),
         "Thank You", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.75), Inches(11.3), Inches(0.6),
         "Questions & Discussion", size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, italic=True)
add_text(s, Inches(1), Inches(5.8), Inches(11.3), Inches(0.4),
         "Gokul Sankar  ·  Banking Data Platform  ·  Batch & Real-Time Streaming Architecture",
         size=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)


prs.save("Banking_Data_Platform_Presentation.pptx")
print("Presentation saved.")



